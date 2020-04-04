import copy

import six
from pptx import Presentation

from tkinter_message import show_error


class MergePPT:
    def __init__(self, matches: dict, presentation, path):
        self.presentation = presentation
        self.matches = matches
        self.save_path = path

    def replace_matches_newslide(self, max_count):
        for i in range(max_count - 1):
            self.duplicate_slide(0)
        for i, items in enumerate(self.matches):
            for key, value in items.items():
                for shapes in self.presentation.slides[i].shapes:
                    if shapes.has_text_frame:
                        if ("{{" + key + "}}") in shapes.text_frame.text:
                            shapes.text_frame.text = shapes.text_frame.text.replace("{{" + key + "}}", str(value))

    def replace_matches(self):
        for key, value in self.matches.items():
            for slide in self.presentation.slides:
                for shapes in slide.shapes:
                    if shapes.has_text_frame:
                        if ("{{" + key + "}}") in shapes.text_frame.text:
                            shapes.text_frame.text = shapes.text_frame.text.replace("{{" + key + "}}", str(value))

    def save(self):
        self.presentation.save(self.save_path)

    def duplicate_slide(self, index):
        pres = self.presentation
        template = pres.slides[index]
        try:
            blank_slide_layout = pres.slide_layouts[12]
        except:
            blank_slide_layout = pres.slide_layouts[len(pres.slide_layouts) - 1]

        copied_slide = pres.slides.add_slide(blank_slide_layout)

        for shp in template.shapes:
            el = shp.element
            newel = copy.deepcopy(el)
            copied_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

        for _, value in six.iteritems(template.part.rels):
            # Make sure we don't copy a notesSlide relation as that won't exist
            if "notesSlide" not in value.reltype:
                copied_slide.part.rels.add_relationship(value.reltype,
                                                        value._target,
                                                        value.rId)

        return copied_slide


def init_ppt(path):
    try:
        return Presentation(path)
    except FileNotFoundError:
        show_error("File not found", "File at " + str(path) + "not found")